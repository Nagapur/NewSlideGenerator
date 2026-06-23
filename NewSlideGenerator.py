!pip install python-pptx
import pandas as pd
from pptx import Presentation
ppt_file_name = 'STE CDCS Promotion Ceremony Deck 2025 v1'
excel_file = '2025-07-10 CDCS-Diff.xlsx'
# Load the Excel file with multiple sheets
xls = pd.ExcelFile(excel_file)


# Load the existing PowerPoint presentation
prs = Presentation(ppt_file_name + ".pptx")
# Function to get a slide layout by name
def get_layout_by_name(prs, layout_name):
    for layout in prs.slide_master.slide_layouts:
        if layout.name == layout_name:
            return layout
    return None
# Function to print placeholder details, to know what their names are
def print_placeholder_details(placeholders):
    for placeholder in placeholders:
        print(f"Index: {placeholder.placeholder_format.idx}, Name: {placeholder
        .name}, Type: {type(placeholder)}")

#get the grade count by a sheet name
def get_count_by_grade(grade):
    dfs = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')

    count = 0

    for sheet_name, df in dfs.items():
        if sheet_name!= 'MSS_Absent' and 'Grade' in df.columns:
            for value in df['Grade']:
                if value == grade:
                    count += 1

    print(f"{grade} Count: {count}")
    return count

def add_E_title_slide():
  e_slide_layout = get_layout_by_name(prs, "MSS_E_Cover")
  slide = prs.slides.add_slide(e_slide_layout)

def add_M_title_slide():
  m_slide_layout = get_layout_by_name(prs, "MSS_M_Cover")
  slide = prs.slides.add_slide(m_slide_layout)

#Create slides by section and grade
def create_slides():

    # Read all sheets once
    dfs = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')

    # Exclude MSS_Absent and combine
    combined_df = pd.concat(
        {k: v for k, v in dfs.items() if k != 'MSS_Absent'},
        names=['SheetName']
    )

    # Define sorting order
    grade_order = ['E', 'M']
    dept_order = ['Corporate', 'Cloud', 'DCS', 'Defence Enterprise']

    # Convert to categorical
    combined_df['Grade'] = pd.Categorical(combined_df['Grade'], categories=grade_order, ordered=True)
    combined_df['BU'] = pd.Categorical(combined_df['BU'], categories=dept_order, ordered=True)

    # Sort
    combined_df = combined_df.sort_values(by=['Grade', 'BU'])

    # Validate required columns
    expected_columns = ['Name', 'Title', 'BU']
    if not all(col in combined_df.columns for col in expected_columns):
        raise ValueError(f"Missing required columns: {expected_columns}")

    prev_grade = None

    # Loop through rows
    for row in combined_df.itertuples():
        sheet_name = row.Index[0]

        # Insert M section slide when switching from E → M
        if prev_grade == 'E' and row.Grade == 'M':
          add_M_title_slide()


        # Get layout for current sheet
        slide_layout = get_layout_by_name(prs, sheet_name)
        if slide_layout is None:
            raise ValueError(f"No layout found for {sheet_name}")

        # Create slide
        slide = prs.slides.add_slide(slide_layout)

        # Fill placeholders
        slide.placeholders[10].text = row.Name
        slide.placeholders[11].text = row.Title
        slide.placeholders[12].text = row.BU

        prev_grade = row.Grade

def add_absent_title_slide():
  a_slide_layout = get_layout_by_name(prs, "MSS_Absent1")
  slide = prs.slides.add_slide(a_slide_layout)

def create_absent_slides(sheet_name):
       # Read all sheets once
    dfs = pd.read_excel(excel_file, sheet_name=sheet_name, engine='openpyxl')

    # Define sorting order
    grade_order = ['E', 'M']
    dept_order = ['Corporate', 'Cloud', 'DCS', 'Defence Enterprise']

    # Convert to categorical
    dfs['Grade'] = pd.Categorical(dfs['Grade'], categories=grade_order, ordered=True)
    dfs['BU'] = pd.Categorical(dfs['BU'], categories=dept_order, ordered=True)

    # Sort
    dfs = dfs.sort_values(by=['Grade', 'BU'])

    # Validate required columns
    expected_columns = ['Name', 'Title', 'BU']
    if not all(col in dfs.columns for col in expected_columns):
        raise ValueError(f"Missing required columns: {expected_columns}")

    prev_grade = None

    # Loop through rows
    for row in dfs.itertuples():


        # Insert M section slide when switching from E → M
        if prev_grade == 'E' and row.Grade == 'M':
            m_slide_layout = get_layout_by_name(prs, "MSS_M_Cover")
            prs.slides.add_slide(m_slide_layout)

        # Get layout for current sheet
        slide_layout = get_layout_by_name(prs, sheet_name)
        if slide_layout is None:
            raise ValueError(f"No layout found for {sheet_name}")

        # Create slide
        slide = prs.slides.add_slide(slide_layout)

        # Fill placeholders
        slide.placeholders[10].text = row.Name
        slide.placeholders[11].text = row.Title
        slide.placeholders[12].text = row.BU

        prev_grade = row.Grade


add_E_title_slide()
create_slides()
add_absent_title_slide()
add_E_title_slide()
create_absent_slides('MSS_Absent')


# Save the new presentation
prs.save(ppt_file_name + '_NEW.pptx')
